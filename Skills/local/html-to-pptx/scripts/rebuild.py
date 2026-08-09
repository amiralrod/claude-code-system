#!/usr/bin/env python3
"""rebuild.py <slides.json> <out.pptx> — rebuilds captured slides as native
PowerPoint objects. Coordinates in slides.json are px on a 1920x1080 frame:
inches = px/144, points = px/2."""
import json
import math
import os
import sys

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

PX_PER_IN = 144.0
SLIDE_W_EMU, SLIDE_H_EMU = 12192000, 6858000
SLIDE_W_PX = 1920.0  # normalised coordinate space width

WEIGHT_SUFFIX = {100: " Thin", 200: " ExtraLight", 300: " Light",
                 500: " Medium", 600: " SemiBold", 800: " ExtraBold",
                 900: " Black"}

def font_name(family, weight):
    return family + WEIGHT_SUFFIX.get(weight, "")

def IN(px):
    return Inches(px / PX_PER_IN)

def PTS(px):
    return Pt(round(px / 2.0, 1))

ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
         "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}

# ---------- low-level fill/effect XML helpers ----------

FILL_TAGS = ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill",
             "a:pattFill", "a:grpFill")

def _spPr(shape):
    return shape._element.spPr

def _clear_fill(spPr):
    for tag in FILL_TAGS:
        for e in spPr.findall(qn(tag)):
            spPr.remove(e)

def _insert_fill(spPr, el):
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        ln.addprevious(el)
    else:
        spPr.append(el)

def _srgb(parent, color, alpha=1.0):
    c = etree.SubElement(parent, qn("a:srgbClr"))
    c.set("val", color)
    if alpha < 1:
        a = etree.SubElement(c, qn("a:alpha"))
        a.set("val", str(int(round(alpha * 100000))))
    return c

def set_solid_fill(shape, color, alpha=1.0):
    spPr = _spPr(shape)
    _clear_fill(spPr)
    sf = spPr.makeelement(qn("a:solidFill"), {})
    _srgb(sf, color, alpha)
    _insert_fill(spPr, sf)

def set_gradient_fill(shape, stops, angle_css):
    spPr = _spPr(shape)
    _clear_fill(spPr)
    gf = spPr.makeelement(qn("a:gradFill"), {})
    lst = etree.SubElement(gf, qn("a:gsLst"))
    for s in stops:
        gs = etree.SubElement(lst, qn("a:gs"))
        gs.set("pos", str(int(round(s["pos"] * 100000))))
        _srgb(gs, s["color"], s.get("alpha", 1))
    lin = etree.SubElement(gf, qn("a:lin"))
    lin.set("ang", str(int(((angle_css - 90) % 360) * 60000)))
    lin.set("scaled", "1")
    _insert_fill(spPr, gf)

def set_shadow(shape, sh):
    spPr = _spPr(shape)
    eff = etree.SubElement(spPr, qn("a:effectLst"))
    o = etree.SubElement(eff, qn("a:outerShdw"))
    o.set("blurRad", str(int(IN(sh["blur"]))))
    o.set("dist", str(int(IN(math.hypot(sh["dx"], sh["dy"])))))
    o.set("dir", str(int((math.degrees(math.atan2(sh["dy"], sh["dx"])) % 360) * 60000)))
    o.set("rotWithShape", "0")
    _srgb(o, sh["color"], sh.get("alpha", 0.3))

def _stroke_alpha(shape, alpha):
    """Add transparency to a shape's outline color (python-pptx has no API)."""
    if alpha >= 1:
        return
    ln = _spPr(shape).find(qn("a:ln"))
    if ln is None:
        return
    sf = ln.find(qn("a:solidFill"))
    c = sf.find(qn("a:srgbClr")) if sf is not None else None
    if c is not None:
        a = etree.SubElement(c, qn("a:alpha"))
        a.set("val", str(int(round(alpha * 100000))))

def _run_alpha(run, alpha):
    rPr = run._r.get_or_add_rPr()
    sf = rPr.find(qn("a:solidFill"))
    if sf is None:
        return
    c = sf.find(qn("a:srgbClr"))
    if c is not None:
        a = etree.SubElement(c, qn("a:alpha"))
        a.set("val", str(int(round(alpha * 100000))))

def _bullet(p, b):
    pPr = p._p.get_or_add_pPr()
    mar = int(IN(b.get("indentPx", 24)))  # hanging indent from the deck's own list padding
    pPr.set("marL", str(mar))
    pPr.set("indent", str(-mar))
    buClr = etree.SubElement(pPr, qn("a:buClr"))
    _srgb(buClr, b["color"])
    buChar = etree.SubElement(pPr, qn("a:buChar"))
    buChar.set("char", b.get("char", "•"))

# ---------- element renderers ----------

def _strip_style(shape):
    """Drop the theme <p:style> (fill/line/effect refs) that add_shape and
    add_connector attach — it injects theme shadows and colors we replace."""
    st = shape._element.find(qn("p:style"))
    if st is not None:
        shape._element.remove(st)

def add_box(slide, el):
    kind = {"ellipse": MSO_SHAPE.OVAL, "roundRect": MSO_SHAPE.ROUNDED_RECTANGLE}
    shape_type = kind.get(el["shape"], MSO_SHAPE.RECTANGLE)
    rotation = 0
    corners = el.get("corners")
    top2 = corners and corners["tl"] and corners["tr"] and not (corners["bl"] or corners["br"])
    bot2 = corners and corners["bl"] and corners["br"] and not (corners["tl"] or corners["tr"])
    if top2 or bot2:
        # partially clipped child: round only the corners shared with the
        # rounded overflow-hidden ancestor (bottom pair = same shape, flipped)
        shape_type = MSO_SHAPE.ROUND_2_SAME_RECTANGLE
        rotation = 180 if bot2 else 0
    elif corners and any(corners.values()):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE  # mixed combos: round all
    frac = min(0.5, el.get("radius", 0) / max(1.0, min(el["w"], el["h"])))
    sp = slide.shapes.add_shape(shape_type,
                                IN(el["x"]), IN(el["y"]), IN(el["w"]), IN(el["h"]))
    _strip_style(sp)
    if rotation:
        sp.rotation = rotation
    if shape_type == MSO_SHAPE.ROUND_2_SAME_RECTANGLE:
        sp.adjustments[0] = frac
        sp.adjustments[1] = 0
    elif el["shape"] == "roundRect" or shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        sp.adjustments[0] = frac
    f = el.get("fill")
    if not f:
        sp.fill.background()
    elif f["type"] == "gradient":
        set_gradient_fill(sp, f["stops"], f["angle"])
    else:
        set_solid_fill(sp, f["color"], f.get("alpha", 1))
    b = el.get("border")
    if b:
        sp.line.color.rgb = RGBColor.from_string(b["color"])
        sp.line.width = PTS(b["width"])
        _stroke_alpha(sp, b.get("alpha", 1))
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    if el.get("shadow"):
        set_shadow(sp, el["shadow"])

def add_line(slide, el):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    IN(el["x1"]), IN(el["y1"]), IN(el["x2"]), IN(el["y2"]))
    _strip_style(ln)
    ln.line.color.rgb = RGBColor.from_string(el["color"])
    ln.line.width = max(Emu(3175), PTS(el["width"]))  # >= 0.25pt
    _stroke_alpha(ln, el.get("alpha", 1))
    ln.shadow.inherit = False

def add_text(slide, el):
    # a block the browser rendered on one line must never wrap in PowerPoint
    single_line = (len(el["paragraphs"]) == 1 and el.get("lineHeightPx")
                   and el["h"] <= el["lineHeightPx"] * 1.6)
    # wrap tolerance: PowerPoint measures text slightly wider than the browser;
    # letter-spaced text needs proportionally more room
    frac = 0.18 if single_line else 0.03
    r0 = el["paragraphs"][0]["runs"][0] if el["paragraphs"][0]["runs"] else {}
    if single_line and r0.get("letterSpacingPx", 0) > 0 and r0.get("fontSizePx"):
        frac += 2.5 * r0["letterSpacingPx"] / r0["fontSizePx"]
    # cap padding so the box never extends past the slide's right edge
    max_pad = max(0.0, SLIDE_W_PX - el["x"] - el["w"])
    pad = min(el["w"] * frac, max_pad)
    x = el["x"]
    if el.get("align") == "center":
        x -= pad / 2
    elif el.get("align") == "right":
        x -= pad
    tb = slide.shapes.add_textbox(IN(x), IN(el["y"]), IN(el["w"] + pad), IN(el["h"]))
    tf = tb.text_frame
    tf.word_wrap = not single_line
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    for i, para in enumerate(el["paragraphs"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ALIGN.get(el.get("align", "left"), PP_ALIGN.LEFT)
        if el.get("lineHeightPx"):
            p.line_spacing = PTS(el["lineHeightPx"])
        if para.get("bullet"):
            _bullet(p, para["bullet"])
        for run in para["runs"]:
            r = p.add_run()
            r.text = run["text"]
            f = r.font
            f.size = PTS(run["fontSizePx"])
            f.name = font_name(run["fontFamily"], run["weight"])
            f.bold = run["weight"] == 700
            f.italic = bool(run.get("italic"))
            f.color.rgb = RGBColor.from_string(run["color"])
            if run.get("alpha", 1) < 1:
                _run_alpha(r, run["alpha"])
            if run.get("letterSpacingPx"):
                r._r.get_or_add_rPr().set(
                    "spc", str(int(run["letterSpacingPx"] / 2 * 100)))

def add_picture(slide, el, asset_dir):
    slide.shapes.add_picture(os.path.join(asset_dir, el["src"]),
                             IN(el["x"]), IN(el["y"]), IN(el["w"]), IN(el["h"]))

def maybe_background(slide, elements):
    """A solid box covering the whole frame becomes the slide background."""
    if elements:
        e = elements[0]
        if (e["type"] == "box" and e["x"] <= 2 and e["y"] <= 2
                and e["w"] >= 1910 and e["h"] >= 1070
                and e.get("fill") and e["fill"]["type"] == "solid"
                and e["fill"].get("alpha", 1) >= 0.99 and not e.get("border")):
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor.from_string(e["fill"]["color"])
            return elements[1:]
    return elements

def build(data, out_path, asset_dir):
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)
    blank = prs.slide_layouts[6]
    problems = []
    for s in data["slides"]:
        slide = prs.slides.add_slide(blank)
        if s.get("failed"):
            if s.get("screenshot"):
                slide.shapes.add_picture(os.path.join(asset_dir, s["screenshot"]),
                                         0, 0, prs.slide_width, prs.slide_height)
            problems.append(f"Slide {s['index'] + 1}: inserted as a picture "
                            "(could not be converted to editable objects).")
            continue
        for el in maybe_background(slide, s["elements"]):
            try:
                if el["type"] == "box":
                    add_box(slide, el)
                elif el["type"] == "line":
                    add_line(slide, el)
                elif el["type"] == "text":
                    add_text(slide, el)
                elif el["type"] in ("image", "svgPicture"):
                    add_picture(slide, el, asset_dir)
            except Exception as exc:  # keep going; report at the end
                problems.append(f"Slide {s['index'] + 1}: skipped one "
                                f"{el['type']} element ({exc}).")
    prs.save(out_path)
    return problems

def main():
    if len(sys.argv) != 3:
        print("Usage: rebuild.py <slides.json> <out.pptx>", file=sys.stderr)
        sys.exit(2)
    json_path, out_path = sys.argv[1], sys.argv[2]
    with open(json_path) as f:
        data = json.load(f)
    problems = build(data, out_path, os.path.dirname(os.path.abspath(json_path)))
    print(f"Wrote {out_path} ({len(data['slides'])} slides)")
    for p in problems + data.get("warnings", []):
        print("WARNING: " + p)

if __name__ == "__main__":
    main()
